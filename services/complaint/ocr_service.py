"""
起诉状模块 OCR 服务
通过 get_ocr_engine() 工厂函数获取当前配置的 OCR 引擎
支持 docx/xlsx/pptx/pdf/图片 的文本提取
"""
from __future__ import annotations

import tempfile
from pathlib import Path
from typing import Optional, Any

from loguru import logger

from config.settings import settings
from services.ocr.base import BaseOCREngine
from services.ocr.engine import get_ocr_engine


_complaint_ocr_engine: Optional[BaseOCREngine] = None


def _get_engine() -> BaseOCREngine:
    """获取当前配置的 OCR 引擎（工厂函数）"""
    global _complaint_ocr_engine
    if _complaint_ocr_engine is None:
        _complaint_ocr_engine = get_ocr_engine()
    return _complaint_ocr_engine


def _extract_docx_text(file_bytes: bytes) -> dict:
    from docx import Document
    import io

    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    tables_text = []
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells]
            tables_text.append(" | ".join(row_text))

    full_text = "\n".join(paragraphs)
    if tables_text:
        full_text += "\n\n[表格内容]\n" + "\n".join(tables_text)

    return {
        "full_text": full_text,
        "blocks": [{"text": p, "source": "paragraph"} for p in paragraphs],
        "block_count": len(paragraphs),
        "source_type": "docx",
    }


def _extract_xlsx_text(file_bytes: bytes) -> dict:
    from openpyxl import load_workbook
    import io

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    all_rows = []
    blocks = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        all_rows.append(f"[工作表: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            row_text = [str(cell) if cell is not None else "" for cell in row]
            line = " | ".join(row_text)
            if line.strip().replace("|", "").strip():
                all_rows.append(line)
                blocks.append({"text": line, "source": f"sheet:{sheet_name}"})

    wb.close()

    return {
        "full_text": "\n".join(all_rows),
        "blocks": blocks,
        "block_count": len(blocks),
        "source_type": "xlsx",
    }


def _extract_pptx_text(file_bytes: bytes) -> dict:
    from pptx import Presentation
    import io

    prs = Presentation(io.BytesIO(file_bytes))
    all_text = []
    blocks = []

    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(row_text))
        if slide_texts:
            all_text.append(f"[幻灯片 {i+1}]")
            all_text.extend(slide_texts)
            for t in slide_texts:
                blocks.append({"text": t, "source": f"slide:{i+1}"})

    return {
        "full_text": "\n".join(all_text),
        "blocks": blocks,
        "block_count": len(blocks),
        "source_type": "pptx",
    }


def ocr_image(image_bytes: bytes, filename: str = "upload.png") -> dict:
    """OCR 识别图片文件

    使用 TemporaryDirectory 确保临时文件在退出时自动清理，
    避免 /tmp 残留撑满 tmpfs。
    """
    engine = _get_engine()

    suffix = Path(filename).suffix.lower() or ".png"
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_path = Path(tmp_dir) / f"input{suffix}"
        tmp_path.write_bytes(image_bytes)

        results = engine.recognize(tmp_path)

    text_lines = []
    for r in results:
        text_lines.append(r.get("text", ""))
    full_text = "\n".join(text_lines)
    return {
        "full_text": full_text,
        "blocks": results,
        "block_count": len(results),
        "source_type": "image_ocr",
    }


def _empty_pdf_result() -> dict:
    return {"full_text": "", "blocks": [], "block_count": 0, "page_count": 0, "source_type": "pdf_ocr"}


def ocr_pdf_path(
    pdf_path: str | Path,
    filename: str = "upload.pdf",
    store: Any | None = None,
) -> dict:
    """OCR 识别磁盘上的 PDF 文件（内存恒定，支持 3000 页超大 PDF）

    关键点：
    - 不把整份 PDF 读进内存（fitz 直接打开磁盘文件，且不做二次拷贝）。
    - 分批：拆一批页 → OCR 一批 → 删一批图，磁盘/内存峰值恒定在单批范围，
      与总页数无关。
    - 临时图片放在磁盘卷(OCR_WORK_DIR)而非内存型 tmpfs(/tmp)。
    """
    import fitz as _fitz
    import os as _os

    try:
        from services.preprocessor.pdf_splitter import PDFSplitter
        from services.ocr.batch_processor import OCRBatchProcessor

        pdf_path = Path(pdf_path)

        use_cloud_ocr = settings.ocr_engine_type in ("bailian", "dashscope", "qwen")
        split_dpi = 150 if use_cloud_ocr else settings.preprocess_dpi

        # 临时目录基路径：优先用磁盘卷(OCR_WORK_DIR)，避免大 PDF 撑满内存型 tmpfs(/tmp)
        work_base = _os.getenv("OCR_WORK_DIR") or None
        if work_base:
            Path(work_base).mkdir(parents=True, exist_ok=True)

        # 每批处理的页数（拆一批→OCR一批→删一批），把磁盘/内存峰值限制在单批范围
        batch_pages = max(1, settings.ocr_batch_size)

        all_text_parts: list[str] = []
        all_blocks: list[dict] = []
        page_count = 0
        block_count = 0
        use_store = store is not None

        # 仅读取页数，不把整份 PDF 载入内存
        doc = _fitz.open(str(pdf_path))
        total_pages = doc.page_count
        doc.close()

        if total_pages == 0:
            return _empty_pdf_result()

        with tempfile.TemporaryDirectory(dir=work_base) as work_dir:
            work_path = Path(work_dir)
            pages_dir = work_path / "pages"
            pages_dir.mkdir()
            ocr_dir = pages_dir / "ocr"

            splitter = PDFSplitter(dpi=split_dpi)
            processor = OCRBatchProcessor()

            # 分批：拆页 → OCR → 删图，循环直到处理完所有页
            for start in range(1, total_pages + 1, batch_pages):
                end = min(start + batch_pages - 1, total_pages)
                image_paths = splitter.split_to_images(
                    pdf_path, pages_dir, start_page=start, end_page=end
                )
                if not image_paths:
                    continue

                batch_summary = processor.process_pages(
                    image_paths, ocr_dir, page_offset=start - 1
                )
                page_count += batch_summary.get("total_pages", len(image_paths))

                for page_data in batch_summary.get("pages", []):
                    page_num = page_data.get("page", 0)
                    results = page_data.get("results", [])
                    if use_store:
                        store.write_page(page_num, results)
                        block_count += len(results)
                    else:
                        for r in results:
                            all_blocks.append({
                                "text": r.get("text", ""),
                                "confidence": r.get("confidence", 0),
                                "page": page_num,
                            })
                            all_text_parts.append(r.get("text", ""))

                # 释放本批图片，保证磁盘/tmpfs 占用恒定在单批范围
                for p in image_paths:
                    try:
                        p.unlink()
                    except OSError:
                        pass

        if page_count == 0:
            return _empty_pdf_result()

        if use_store:
            return {
                "full_text": "",
                "blocks": [],
                "block_count": block_count,
                "page_count": page_count,
                "source_type": "pdf_ocr",
                "offloaded": True,
            }

        return {
            "full_text": "\n".join(all_text_parts),
            "blocks": all_blocks,
            "block_count": len(all_blocks),
            "page_count": page_count,
            "source_type": "pdf_ocr",
        }
    except Exception:
        logger.exception(f"ocr_pdf_path failed for {filename}")
        raise


def ocr_pdf_page_range(
    pdf_path: str | Path,
    filename: str,
    start_page: int,
    end_page: int,
    store: Any | None = None,
) -> dict:
    """OCR 识别磁盘 PDF 的指定页段 [start_page, end_page]（1-based 闭区间）。

    专供分片批次 task 使用：fitz 打开磁盘文件，只渲染 [start_page, end_page]，
    OCR 结果逐页写入 store（EvidenceOCRStore），不拼接 full_text（收口器统一拼）。

    内存恒定：内部仍按 ocr_batch_size 拆子批，拆一批→OCR→删一批。
    """
    import os as _os
    from pathlib import Path as _Path

    try:
        from services.preprocessor.pdf_splitter import PDFSplitter
        from services.ocr.batch_processor import OCRBatchProcessor

        pdf_path = _Path(pdf_path)
        use_cloud_ocr = settings.ocr_engine_type in ("bailian", "dashscope", "qwen")
        split_dpi = 150 if use_cloud_ocr else settings.preprocess_dpi

        work_base = _os.getenv("OCR_WORK_DIR") or None
        if work_base:
            _Path(work_base).mkdir(parents=True, exist_ok=True)

        # 子批大小（仍按全局 ocr_batch_size 拆，内存峰值恒定）
        sub_batch = max(1, settings.ocr_batch_size)
        block_count = 0
        page_count = 0

        with tempfile.TemporaryDirectory(dir=work_base) as work_dir:
            work_path = _Path(work_dir)
            pages_dir = work_path / "pages"
            pages_dir.mkdir()
            ocr_dir = pages_dir / "ocr"

            splitter = PDFSplitter(dpi=split_dpi)
            processor = OCRBatchProcessor()

            for sub_start in range(start_page, end_page + 1, sub_batch):
                sub_end = min(sub_start + sub_batch - 1, end_page)
                image_paths = splitter.split_to_images(
                    pdf_path, pages_dir, start_page=sub_start, end_page=sub_end
                )
                if not image_paths:
                    continue

                batch_summary = processor.process_pages(
                    image_paths, ocr_dir, page_offset=sub_start - 1
                )
                page_count += batch_summary.get("total_pages", len(image_paths))

                for page_data in batch_summary.get("pages", []):
                    page_num = page_data.get("page", 0)
                    results = page_data.get("results", [])
                    if store is not None:
                        store.write_page(page_num, results)
                        block_count += len(results)

                for p in image_paths:
                    try:
                        p.unlink()
                    except OSError:
                        pass

        return {
            "full_text": "",
            "blocks": [],
            "block_count": block_count,
            "page_count": page_count,
            "start_page": start_page,
            "end_page": end_page,
            "source_type": "pdf_ocr_shard_batch",
            "offloaded": store is not None,
        }
    except Exception:
        logger.exception(
            f"ocr_pdf_page_range failed for {filename} [{start_page}-{end_page}]"
        )
        raise


def ocr_pdf(pdf_bytes: bytes, filename: str = "upload.pdf") -> dict:
    """OCR 识别 PDF（字节入口，向后兼容）

    大文件应优先用 ocr_pdf_path 走磁盘路径，避免把整份 PDF 读进内存。
    这里把字节先落到磁盘临时文件，再委托给 ocr_pdf_path。
    """
    import os as _os

    work_base = _os.getenv("OCR_WORK_DIR") or None
    if work_base:
        Path(work_base).mkdir(parents=True, exist_ok=True)

    with tempfile.NamedTemporaryFile(dir=work_base, suffix=".pdf", delete=False) as tmp:
        tmp.write(pdf_bytes)
        tmp_path = tmp.name
    try:
        return ocr_pdf_path(tmp_path, filename)
    finally:
        try:
            Path(tmp_path).unlink()
        except OSError:
            pass


OFFICE_SUFFIXES = {".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt"}


def ocr_upload(file_bytes: bytes, filename: str) -> dict:
    suffix = Path(filename).suffix.lower()

    if suffix == ".pdf":
        return ocr_pdf(file_bytes, filename)

    if suffix in (".docx", ".doc"):
        return _extract_docx_text(file_bytes)

    if suffix in (".xlsx", ".xls"):
        return _extract_xlsx_text(file_bytes)

    if suffix in (".pptx", ".ppt"):
        return _extract_pptx_text(file_bytes)

    return ocr_image(file_bytes, filename)


def ocr_upload_path(
    file_path: str | Path,
    filename: str,
    store: Any | None = None,
) -> dict:
    """从磁盘文件入口做 OCR（内存恒定，推荐用于大文件）

    PDF 走 ocr_pdf_path（全程流式，不进内存）；其余格式（docx/xlsx/pptx/图片）
    通常较小，从磁盘读取字节再走原有提取逻辑。
    """
    suffix = Path(filename).suffix.lower()
    file_path = Path(file_path)

    if suffix == ".pdf":
        return ocr_pdf_path(file_path, filename, store=store)

    if suffix in (".docx", ".doc"):
        return _extract_docx_text(file_path.read_bytes())

    if suffix in (".xlsx", ".xls"):
        return _extract_xlsx_text(file_path.read_bytes())

    if suffix in (".pptx", ".ppt"):
        return _extract_pptx_text(file_path.read_bytes())

    result = ocr_image(file_path.read_bytes(), filename)
    if store is not None:
        store.set_meta("image_ocr")
        store.write_page(1, result.get("blocks") or [])
        return {
            "full_text": "",
            "blocks": [],
            "block_count": result.get("block_count", 0),
            "page_count": 1,
            "source_type": "image_ocr",
            "offloaded": True,
        }
    return result
